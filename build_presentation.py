"""Build final_presentation.pptx for the AutoResearch Spotify project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
import os

# ── colour palette ──────────────────────────────────────────────────────────
SPOTIFY_GREEN  = RGBColor(0x1D, 0xB9, 0x54)
DARK_BG        = RGBColor(0x12, 0x12, 0x12)
LIGHT_BG       = RGBColor(0xF4, 0xF4, 0xF4)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT      = RGBColor(0x1A, 0x1A, 0x1A)
MID_GRAY       = RGBColor(0x55, 0x55, 0x55)
ACCENT_RED     = RGBColor(0xE0, 0x3A, 0x3A)
ACCENT_AMBER   = RGBColor(0xF5, 0xA6, 0x23)
ACCENT_GREEN   = RGBColor(0x27, 0xAE, 0x60)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.line.width = 0
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_rgb
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline_text(slide, lines, l, t, w, h,
                       font_size=16, color=DARK_TEXT, bold=False,
                       align=PP_ALIGN.LEFT, line_spacing_pt=None,
                       font_name="Calibri"):
    """lines is a list of (text, bold, color, size) tuples or plain strings."""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, b, c, sz = line, bold, color, font_size
        else:
            txt, b, c, sz = line
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.bold = b
        run.font.color.rgb = c
        run.font.name = font_name
    return txBox


def add_image(slide, path, l, t, w, h=None):
    if h is None:
        slide.shapes.add_picture(path, l, t, width=w)
    else:
        slide.shapes.add_picture(path, l, t, width=w, height=h)


def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


# ── slide background helper ──────────────────────────────────────────────────

def dark_slide(title_text, subtitle_text=None):
    """Returns a slide with dark header bar and light body."""
    slide = prs.slides.add_slide(blank_layout)
    # full background — light
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_BG)
    # dark header bar
    header_h = Inches(1.35)
    add_rect(slide, 0, 0, SLIDE_W, header_h, DARK_BG)
    # green accent line under header
    add_rect(slide, 0, header_h, SLIDE_W, Inches(0.055), SPOTIFY_GREEN)
    # title text
    add_text(slide, title_text,
             Inches(0.45), Inches(0.28), Inches(12.5), Inches(0.85),
             font_size=32, bold=True, color=WHITE)
    if subtitle_text:
        add_text(slide, subtitle_text,
                 Inches(0.45), Inches(0.95), Inches(12.5), Inches(0.45),
                 font_size=15, color=SPOTIFY_GREEN)
    return slide


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(blank_layout)
add_rect(slide1, 0, 0, SLIDE_W, SLIDE_H, DARK_BG)
# big green accent stripe
add_rect(slide1, 0, Inches(3.3), SLIDE_W, Inches(0.08), SPOTIFY_GREEN)

add_text(slide1,
         "The Limits of Audio-Only Popularity Prediction",
         Inches(0.6), Inches(1.1), Inches(12.2), Inches(1.8),
         font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide1,
         "An AutoResearch Investigation of Spotify Track Popularity",
         Inches(0.6), Inches(2.85), Inches(12.2), Inches(0.6),
         font_size=21, bold=False, color=SPOTIFY_GREEN, align=PP_ALIGN.CENTER)

add_text(slide1,
         "Emily Yoo  ·  AutoResearch Program  ·  Weeks 2–6",
         Inches(0.6), Inches(3.55), Inches(12.2), Inches(0.5),
         font_size=16, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)

# three key-number boxes
box_w, box_h = Inches(3.2), Inches(1.35)
for i, (num, label) in enumerate([
        ("~114,000", "Spotify tracks"),
        ("16", "experiments over 5 weeks"),
        ("22.06 → 15.72", "RMSE baseline → final"),
]):
    bx = Inches(0.55) + i * Inches(4.26)
    by = Inches(4.5)
    add_rect(slide1, bx, by, box_w, box_h, RGBColor(0x22, 0x22, 0x22))
    add_rect(slide1, bx, by + box_h - Inches(0.06), box_w, Inches(0.06), SPOTIFY_GREEN)
    add_text(slide1, num,
             bx + Inches(0.15), by + Inches(0.1), box_w - Inches(0.3), Inches(0.7),
             font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide1, label,
             bx + Inches(0.15), by + Inches(0.75), box_w - Inches(0.3), Inches(0.55),
             font_size=13, color=RGBColor(0xBB, 0xBB, 0xBB), align=PP_ALIGN.CENTER)

add_speaker_notes(slide1,
    "My project asked whether numeric audio features — things like energy, danceability, and "
    "tempo — can predict how popular a song becomes on Spotify. The short answer is: better "
    "than a linear model, but with a hard ceiling that no amount of modeling can break through. "
    "That ceiling is the real finding.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Topic → AutoResearch Contract
# ════════════════════════════════════════════════════════════════════════════
slide2 = dark_slide("The AutoResearch Contract",
                    "Topic → how the research loop was structured")

# two-column layout
col_l = Inches(0.45)
col_mid = Inches(6.8)
body_t = Inches(1.6)
body_h = Inches(5.5)

# left box — research question
add_rect(slide2, col_l, body_t, Inches(6.0), body_h, WHITE)
add_rect(slide2, col_l, body_t, Inches(6.0), Inches(0.5), SPOTIFY_GREEN)
add_text(slide2, "Research Question",
         col_l + Inches(0.15), body_t + Inches(0.06), Inches(5.7), Inches(0.4),
         font_size=14, bold=True, color=WHITE)

rq_lines = [
    "Can Spotify track popularity be predicted",
    "from numeric audio features alone?",
    "",
    "13 numeric features:",
    "danceability · energy · loudness · tempo",
    "speechiness · acousticness · instrumentalness",
    "liveness · valence · key · mode",
    "time_signature · duration_ms",
    "",
    "Target: popularity  (0–100 integer)",
    "Dataset: 114,000 tracks",
    "Evaluation: validation RMSE",
    "Random seed: 42 (frozen throughout)",
]
add_multiline_text(slide2, rq_lines,
                   col_l + Inches(0.2), body_t + Inches(0.65),
                   Inches(5.6), Inches(4.7),
                   font_size=15, color=DARK_TEXT)

# right box — loop contract
add_rect(slide2, col_mid, body_t, Inches(6.1), body_h, WHITE)
add_rect(slide2, col_mid, body_t, Inches(6.1), Inches(0.5), DARK_BG)
add_text(slide2, "The Keep / Discard / Crash Loop",
         col_mid + Inches(0.15), body_t + Inches(0.06), Inches(5.8), Inches(0.4),
         font_size=14, bold=True, color=SPOTIFY_GREEN)

loop_lines = [
    ("Run experiment → measure validation RMSE", False, DARK_TEXT, 15),
    ("", False, DARK_TEXT, 8),
    ("✓  Keep    — RMSE improves → retain change", True, ACCENT_GREEN, 15),
    ("✗  Discard — RMSE worsens → revert", True, ACCENT_RED, 15),
    ("⚡  Crash   — code fails → revert + log", True, ACCENT_AMBER, 15),
    ("", False, DARK_TEXT, 8),
    ("Constraints:", True, DARK_TEXT, 14),
    ("• Frozen train/validation/test split (60/20/20)", False, MID_GRAY, 14),
    ("• Fixed random_state = 42 throughout", False, MID_GRAY, 14),
    ("• No external data, no neural architectures", False, MID_GRAY, 14),
    ("• Test set untouched during all 16 experiments", False, MID_GRAY, 14),
    ("• Evaluation logic in run.py is frozen", False, MID_GRAY, 14),
    ("", False, DARK_TEXT, 8),
    ("16 experiments · 5 active weeks · 0 crashes", True, DARK_TEXT, 15),
]
add_multiline_text(slide2, loop_lines,
                   col_mid + Inches(0.2), body_t + Inches(0.65),
                   Inches(5.8), Inches(4.7),
                   font_size=15, color=DARK_TEXT)

add_speaker_notes(slide2,
    "The AutoResearch loop had a simple rule: run an experiment, measure validation RMSE, "
    "keep it if it improves, discard it if it doesn't. The loop is fully reproducible — "
    "frozen train/validation split, fixed random seed, no external data. That discipline is "
    "what makes the 16 RMSE numbers comparable to each other.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Directions Explored (with metric trajectory chart)
# ════════════════════════════════════════════════════════════════════════════
slide3 = dark_slide("Directions Explored",
                    "What the loop tried across 16 experiments")

# chart takes the right ~60% of the slide
img_path = "week5/metric_trajectory.png"
img_l = Inches(5.4)
img_t = Inches(1.5)
img_w = Inches(7.6)
add_image(slide3, img_path, img_l, img_t, img_w)

# table on left
tbl_l = Inches(0.4)
tbl_t = Inches(1.6)

directions = [
    ("RF adoption",          "Major gain",    ACCENT_GREEN),
    ("Hyperparameter tuning","Diminishing",   ACCENT_AMBER),
    ("Feature engineering",  "Minimal gain",  ACCENT_AMBER),
    ("Alt. model families",  "Discarded",     ACCENT_RED),
    ("Popularity bucket",    "Key insight",   SPOTIFY_GREEN),
    ("Feature bucket",       "Key insight",   SPOTIFY_GREEN),
    ("Index artifact",       "Key insight",   SPOTIFY_GREEN),
]

row_h = Inches(0.52)
row_w = Inches(4.7)
hdr_h = Inches(0.45)

# header
add_rect(slide3, tbl_l, tbl_t, row_w, hdr_h, DARK_BG)
add_text(slide3, "Direction",
         tbl_l + Inches(0.1), tbl_t + Inches(0.05), Inches(2.7), hdr_h,
         font_size=13, bold=True, color=WHITE)
add_text(slide3, "Outcome",
         tbl_l + Inches(2.85), tbl_t + Inches(0.05), Inches(1.8), hdr_h,
         font_size=13, bold=True, color=WHITE)

for i, (direction, outcome, color) in enumerate(directions):
    y = tbl_t + hdr_h + i * row_h
    bg = WHITE if i % 2 == 0 else RGBColor(0xEE, 0xEE, 0xEE)
    add_rect(slide3, tbl_l, y, row_w, row_h, bg)
    # colored outcome badge
    badge_w = Inches(1.55)
    add_rect(slide3, tbl_l + Inches(2.9), y + Inches(0.07),
             badge_w, row_h - Inches(0.14), color)
    add_text(slide3, direction,
             tbl_l + Inches(0.1), y + Inches(0.08), Inches(2.7), row_h - Inches(0.1),
             font_size=13, color=DARK_TEXT)
    add_text(slide3, outcome,
             tbl_l + Inches(2.95), y + Inches(0.08), badge_w - Inches(0.1), row_h - Inches(0.1),
             font_size=12, bold=True, color=WHITE)

add_speaker_notes(slide3,
    "The trajectory chart on the right shows this cleanly. There is one large step down — "
    "the switch to Random Forest — and then the curve flattens. Everything after that first "
    "drop produced combined improvement of 0.06 RMSE units across 12 experiments. The loop "
    "explored the right directions; it just ran out of signal to exploit.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Stable Value (with week4 metric_over_time chart)
# ════════════════════════════════════════════════════════════════════════════
slide4 = dark_slide("Stable Value",
                    "What produced real, reproducible improvement")

# chart — left side
img_path4 = "week4/metric_over_time.png"
img_l4 = Inches(0.35)
img_t4 = Inches(1.5)
img_w4 = Inches(7.0)
add_image(slide4, img_path4, img_l4, img_t4, img_w4)

# stats on right
rx = Inches(7.7)
ry = Inches(1.55)
rw = Inches(5.35)

# headline stat box
add_rect(slide4, rx, ry, rw, Inches(1.5), SPOTIFY_GREEN)
add_text(slide4, "31.9% RMSE reduction",
         rx + Inches(0.15), ry + Inches(0.1), rw - Inches(0.3), Inches(0.7),
         font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide4, "One architectural decision",
         rx + Inches(0.15), ry + Inches(0.8), rw - Inches(0.3), Inches(0.55),
         font_size=17, color=WHITE, align=PP_ALIGN.CENTER)

# comparison rows
cmp_t = ry + Inches(1.65)
for j, (model, rmse, c) in enumerate([
        ("Linear Regression  (baseline)", "RMSE 22.06", ACCENT_RED),
        ("Random Forest  (final model)",  "RMSE 15.72", ACCENT_GREEN),
]):
    yy = cmp_t + j * Inches(1.05)
    add_rect(slide4, rx, yy, rw, Inches(0.95), c)
    add_text(slide4, model,
             rx + Inches(0.2), yy + Inches(0.08), rw * 0.6, Inches(0.4),
             font_size=15, bold=True, color=WHITE)
    add_text(slide4, rmse,
             rx + rw * 0.62, yy + Inches(0.08), rw * 0.35, Inches(0.4),
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

# breakdown bullets
bullet_t = cmp_t + Inches(2.35)
bullets = [
    ("99% of total gain", "from the model-family switch alone"),
    ("0.06 RMSE units", "from all 12 post-RF experiments combined"),
    ("n_estimators 100→300", "6× runtime for 0.5% additional gain"),
]
for j, (stat, desc) in enumerate(bullets):
    yy = bullet_t + j * Inches(0.72)
    add_text(slide4, stat,
             rx, yy, rw * 0.44, Inches(0.6),
             font_size=15, bold=True, color=DARK_TEXT)
    add_text(slide4, desc,
             rx + rw * 0.45, yy, rw * 0.55, Inches(0.6),
             font_size=14, color=MID_GRAY)

add_speaker_notes(slide4,
    "The single most important decision in the entire project was switching from Linear "
    "Regression to Random Forest. That one change accounts for 99% of the total improvement. "
    "Everything else — transforms, hyperparameter tuning, alternative ensembles — moved the "
    "needle by a combined 0.06 units. If the loop had stopped after Week 2, the result would "
    "have been essentially the same.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Noise and Failure
# ════════════════════════════════════════════════════════════════════════════
slide5 = dark_slide("Noise and Failure",
                    "10 of 16 experiments were discarded — and why that matters")

body_t5 = Inches(1.55)
col1_l = Inches(0.45)
col2_l = Inches(6.9)
col_w  = Inches(6.05)
col_h  = Inches(5.6)

# left panel — what was discarded
add_rect(slide5, col1_l, body_t5, col_w, col_h, WHITE)
add_rect(slide5, col1_l, body_t5, col_w, Inches(0.5), ACCENT_RED)
add_text(slide5, "Discarded (10 experiments)",
         col1_l + Inches(0.15), body_t5 + Inches(0.06), col_w - Inches(0.3), Inches(0.42),
         font_size=14, bold=True, color=WHITE)

disc_rows = [
    ("Interaction terms  ×3",       "RF discovers interactions natively",  "+0.004–+0.019"),
    ("HistGradientBoosting (×2)",    "Worse than RF even when tuned",       "+2.73 / +3.38"),
    ("ExtraTreesRegressor",          "Random splits < RF best-split",       "+0.43"),
    ("RF max_features / min_samples","Regularization not the bottleneck",   "+0.26–+0.36"),
    ("Log transforms (bounded feats)","Bounded features — log adds noise",   "+0.003–+0.007"),
    ("Feature dropping",             "key, mode, time_sig still contribute","+0.71"),
]
row_h5 = Inches(0.73)
for i, (name, reason, delta) in enumerate(disc_rows):
    yy = body_t5 + Inches(0.6) + i * row_h5
    bg = WHITE if i % 2 == 0 else RGBColor(0xF8, 0xF0, 0xF0)
    add_rect(slide5, col1_l, yy, col_w, row_h5, bg)
    add_text(slide5, name,
             col1_l + Inches(0.12), yy + Inches(0.04), col_w * 0.42, Inches(0.38),
             font_size=13, bold=True, color=DARK_TEXT)
    add_text(slide5, reason,
             col1_l + Inches(0.12), yy + Inches(0.38), col_w * 0.63, Inches(0.34),
             font_size=11, color=MID_GRAY)
    add_text(slide5, delta,
             col1_l + col_w * 0.68, yy + Inches(0.18), col_w * 0.3, Inches(0.38),
             font_size=12, bold=True, color=ACCENT_RED, align=PP_ALIGN.RIGHT)

# right panel — structural lesson
add_rect(slide5, col2_l, body_t5, col_w, col_h, WHITE)
add_rect(slide5, col2_l, body_t5, col_w, Inches(0.5), DARK_BG)
add_text(slide5, "What This Reveals",
         col2_l + Inches(0.15), body_t5 + Inches(0.06), col_w - Inches(0.3), Inches(0.42),
         font_size=14, bold=True, color=SPOTIFY_GREEN)

lesson_lines = [
    ("Combined effect of all discarded work:", False, DARK_TEXT, 15),
    ("  0.00 RMSE improvement", True, ACCENT_RED, 22),
    ("", False, DARK_TEXT, 8),
    ("Structural limitation of a pure RMSE loop:", True, DARK_TEXT, 14),
    ("", False, DARK_TEXT, 6),
    ("• Loop executes checklist instead of updating", False, MID_GRAY, 14),
    ("  beliefs after first failure", False, MID_GRAY, 14),
    ("", False, DARK_TEXT, 6),
    ("• Three consecutive interaction experiments ran", False, MID_GRAY, 14),
    ("  even after the first failure made the reason", False, MID_GRAY, 14),
    ("  clear: RF already handles interactions", False, MID_GRAY, 14),
    ("", False, DARK_TEXT, 6),
    ("• No self-auditing → index artifact undetected", False, MID_GRAY, 14),
    ("  for 5 weeks (35.8% feature importance!)", False, MID_GRAY, 14),
    ("", False, DARK_TEXT, 8),
    ("Fix: differentiate evidence by replaceability", True, DARK_TEXT, 14),
    ("of hypothesis, not flat keep/discard rule", False, MID_GRAY, 13),
]
add_multiline_text(slide5, lesson_lines,
                   col2_l + Inches(0.2), body_t5 + Inches(0.6),
                   col_w - Inches(0.35), col_h - Inches(0.7),
                   font_size=14, color=DARK_TEXT)

add_speaker_notes(slide5,
    "Documenting failure is as important as documenting success. Three consecutive "
    "interaction-term experiments all failed for the same reason — RF already discovers "
    "interactions through splits — but the loop ran all three before moving on. That's a "
    "structural weakness: the loop executes its checklist rather than updating its beliefs "
    "after the first failure.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Most Important Discovery: Popularity Bucket Analysis
# ════════════════════════════════════════════════════════════════════════════
slide6 = dark_slide("The Model Fails Where It Matters Most",
                    "Popularity bucket analysis — the central result")

# chart
img_path6 = "week6/error_by_bucket_plot.png"
img_l6 = Inches(0.35)
img_t6 = Inches(1.5)
img_w6 = Inches(7.2)
add_image(slide6, img_path6, img_l6, img_t6, img_w6)

# right side
rx6 = Inches(7.9)
ry6 = Inches(1.55)
rw6 = Inches(5.1)

# bucket results table
for j, (bucket, rmse, bias, note, c) in enumerate([
        ("Low  (0–33)",    "15.22", "+10.67", "overpredicts",       ACCENT_AMBER),
        ("Medium  (34–66)","13.31", "−8.79",  "near-balanced",      ACCENT_GREEN),
        ("High  (67–100)", "28.48", "−23.26", "severe underpredict", ACCENT_RED),
]):
    yy = ry6 + j * Inches(1.2)
    add_rect(slide6, rx6, yy, rw6, Inches(1.1), c)
    add_text(slide6, bucket,
             rx6 + Inches(0.15), yy + Inches(0.05), rw6 * 0.48, Inches(0.45),
             font_size=16, bold=True, color=WHITE)
    add_text(slide6, f"RMSE {rmse}",
             rx6 + rw6 * 0.5, yy + Inches(0.05), rw6 * 0.5, Inches(0.45),
             font_size=19, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    add_text(slide6, f"Bias {bias}  ·  {note}",
             rx6 + Inches(0.15), yy + Inches(0.55), rw6 - Inches(0.3), Inches(0.45),
             font_size=13, color=WHITE)

# key insight text
ki_t = ry6 + Inches(3.7)
add_rect(slide6, rx6, ki_t, rw6, Inches(2.4), RGBColor(0x1A, 0x1A, 0x1A))
add_rect(slide6, rx6, ki_t, rw6, Inches(0.07), ACCENT_RED)
add_multiline_text(slide6, [
    ("2.1× higher RMSE on high-popularity tracks", True, WHITE, 15),
    ("", False, WHITE, 7),
    ("Audio features identify what a song sounds like.", False, RGBColor(0xCC, 0xCC, 0xCC), 13),
    ("They cannot identify which good-sounding songs", False, RGBColor(0xCC, 0xCC, 0xCC), 13),
    ("become hits.", False, RGBColor(0xCC, 0xCC, 0xCC), 13),
    ("", False, WHITE, 7),
    ("Missing signal:", True, SPOTIFY_GREEN, 13),
    ("artist reputation · playlist placement", False, RGBColor(0xBB, 0xBB, 0xBB), 13),
    ("label support · release timing", False, RGBColor(0xBB, 0xBB, 0xBB), 13),
],
    rx6 + Inches(0.2), ki_t + Inches(0.2), rw6 - Inches(0.35), Inches(2.15),
    font_size=13, color=WHITE)

add_speaker_notes(slide6,
    "This is the central result. The model performs reasonably on the average track but fails "
    "dramatically on breakout popular tracks. That's a 2.1× RMSE difference between the "
    "medium and high buckets. The high-popularity bucket is where all the interesting "
    "prediction happens — and that's exactly where the audio features run out of signal. "
    "The missing information is artist reputation, playlist placement, and release context. "
    "None of that is in the dataset.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Limits of the Loop (Feature Bucket Analysis)
# ════════════════════════════════════════════════════════════════════════════
slide7 = dark_slide("Limits of the Loop",
                    "Feature-bucket analysis: the failure is not in audio space")

# chart
img_path7 = "week6/feature_bucket_error_plot.png"
img_l7 = Inches(0.35)
img_t7 = Inches(1.5)
img_w7 = Inches(7.3)
add_image(slide7, img_path7, img_l7, img_t7, img_w7)

rx7 = Inches(7.95)
ry7 = Inches(1.55)
rw7 = Inches(5.05)

add_text(slide7, "Error structure by slicing dimension:",
         rx7, ry7, rw7, Inches(0.45),
         font_size=14, bold=True, color=DARK_TEXT)

rows7 = [
    ("Popularity buckets", "15.2 units", "33.9 units", ACCENT_RED),
    ("Danceability",       "3.1 units",  "1.1 units",  ACCENT_GREEN),
    ("Energy",             "0.3 units",  "0.4 units",  ACCENT_GREEN),
    ("Tempo",              "1.4 units",  "0.4 units",  ACCENT_GREEN),
]

# header row
hdr_y = ry7 + Inches(0.5)
add_rect(slide7, rx7, hdr_y, rw7, Inches(0.42), DARK_BG)
for col_x, lbl in [
        (rx7 + Inches(0.1),  "Sliced by"),
        (rx7 + Inches(2.5),  "RMSE span"),
        (rx7 + Inches(3.7),  "Bias span"),
]:
    add_text(slide7, lbl, col_x, hdr_y + Inches(0.04), Inches(1.4), Inches(0.38),
             font_size=12, bold=True, color=WHITE)

row_h7 = Inches(0.62)
for i, (dim, rmse_span, bias_span, c) in enumerate(rows7):
    yy = hdr_y + Inches(0.42) + i * row_h7
    bg = RGBColor(0xFF, 0xEE, 0xEE) if c == ACCENT_RED else RGBColor(0xF0, 0xFF, 0xF0)
    add_rect(slide7, rx7, yy, rw7, row_h7, bg)
    # highlight bar on left
    add_rect(slide7, rx7, yy, Inches(0.07), row_h7, c)
    add_text(slide7, dim,
             rx7 + Inches(0.15), yy + Inches(0.1), Inches(2.2), row_h7 - Inches(0.1),
             font_size=13, bold=(c == ACCENT_RED), color=DARK_TEXT)
    add_text(slide7, rmse_span,
             rx7 + Inches(2.45), yy + Inches(0.1), Inches(1.2), row_h7 - Inches(0.1),
             font_size=13, bold=(c == ACCENT_RED), color=c)
    add_text(slide7, bias_span,
             rx7 + Inches(3.65), yy + Inches(0.1), Inches(1.3), row_h7 - Inches(0.1),
             font_size=13, bold=(c == ACCENT_RED), color=c)

# conclusion box
conc_t = hdr_y + Inches(0.42) + len(rows7) * row_h7 + Inches(0.25)
add_rect(slide7, rx7, conc_t, rw7, Inches(1.85), DARK_BG)
add_rect(slide7, rx7, conc_t, rw7, Inches(0.07), ACCENT_RED)
add_multiline_text(slide7, [
    ("15-unit RMSE gap in popularity space.", True, WHITE, 14),
    ("3-unit gap in audio space.", True, ACCENT_GREEN, 14),
    ("", False, WHITE, 7),
    ("The loop cannot fix what the data", False, RGBColor(0xCC, 0xCC, 0xCC), 13),
    ("does not contain.", False, RGBColor(0xCC, 0xCC, 0xCC), 13),
],
    rx7 + Inches(0.2), conc_t + Inches(0.15), rw7 - Inches(0.35), Inches(1.6),
    font_size=13, color=WHITE)

add_speaker_notes(slide7,
    "This analysis rules out the alternative explanation — that the model fails because it "
    "misreads certain audio profiles. Slicing by any audio feature produces at most 3 units "
    "of RMSE variation and near-zero bias. Slicing by popularity produces 15 units of RMSE "
    "variation and 34 units of bias swing. The failure is in the popularity dimension, not "
    "the audio dimension. No model change, transform, or hyperparameter could have addressed "
    "this. The loop hit the ceiling of its own data.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Reflection and Takeaway
# ════════════════════════════════════════════════════════════════════════════
slide8 = dark_slide("What the Loop Can and Cannot Do",
                    "Reflection · takeaways · future work")

body_t8 = Inches(1.55)

# three-column layout
col_w8 = Inches(3.95)
gap = Inches(0.25)

for ci, (title, items, bg_c, hdr_c) in enumerate([
        ("Loop Strengths", [
            "Zero crashes across 16 experiments",
            "Correct rollbacks on every revert",
            "Identified diminishing returns accurately",
            "Applied log transforms correctly\n(zero-inflated features only)",
            "Reproducible RMSE comparisons",
         ], RGBColor(0xF0, 0xFF, 0xF4), ACCENT_GREEN),
        ("Loop Limits", [
            "Index artifact undetected 5 weeks\n(35.8% importance — no self-audit)",
            "Ran 3 identical interaction experiments\ninstead of stopping after first fail",
            "Analysis started in Week 6 — should\nhave been Week 3",
            "Flat keep/discard rule doesn't weight\nevidence by hypothesis replaceability",
         ], RGBColor(0xFF, 0xF4, 0xF0), ACCENT_RED),
        ("Future Work", [
            "Artist-level features\n(follower count, prior release history)",
            "Playlist placement signals\n(editorial vs. algorithmic)",
            "Release context\n(label, timing, cultural moment)",
            "Social engagement signals",
            "Feature audit checkpoint after\nfirst model fit",
         ], RGBColor(0xF0, 0xF4, 0xFF), SPOTIFY_GREEN),
]):
    cx = Inches(0.35) + ci * (col_w8 + gap)
    add_rect(slide8, cx, body_t8, col_w8, Inches(5.6), bg_c)
    add_rect(slide8, cx, body_t8, col_w8, Inches(0.5), hdr_c)
    add_text(slide8, title,
             cx + Inches(0.12), body_t8 + Inches(0.06), col_w8 - Inches(0.24), Inches(0.42),
             font_size=14, bold=True, color=WHITE)
    for bi, item in enumerate(items):
        yy = body_t8 + Inches(0.62) + bi * Inches(0.92)
        # bullet marker
        add_rect(slide8, cx + Inches(0.12), yy + Inches(0.12),
                 Inches(0.07), Inches(0.07), hdr_c)
        add_text(slide8, item,
                 cx + Inches(0.28), yy, col_w8 - Inches(0.4), Inches(0.85),
                 font_size=13, color=DARK_TEXT)

# final takeaway bar at bottom
bar_t = body_t8 + Inches(5.75)
add_rect(slide8, Inches(0.35), bar_t, SLIDE_W - Inches(0.7), Inches(0.85), DARK_BG)
add_rect(slide8, Inches(0.35), bar_t, Inches(0.09), Inches(0.85), SPOTIFY_GREEN)
add_text(slide8,
         "Improving prediction required understanding the limits of the data, not building more complex models.",
         Inches(0.55), bar_t + Inches(0.12), SLIDE_W - Inches(0.9), Inches(0.65),
         font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_speaker_notes(slide8,
    "The loop is good at mechanical discipline and bad at knowing when to stop or what to "
    "look at. The index artifact — a CSV row number that acted as a genre proxy — was the "
    "highest-importance feature for five weeks because the loop had no self-auditing "
    "component. It was only caught when the interpretability analysis ran in Week 6. The fix "
    "is procedural: inspect feature importances after the first model fit, run a bucket "
    "analysis after Week 3, and build diminishing-returns checkpoints into the loop. The "
    "bigger lesson is that the most valuable output of this project was not the RMSE number "
    "— it was characterizing exactly why the number cannot get lower.")


# ── save ─────────────────────────────────────────────────────────────────────
out_path = "final_presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
